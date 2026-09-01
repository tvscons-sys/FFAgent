"""Load supported source files into LangChain Documents."""

import csv
from pathlib import Path

from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_core.documents import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".csv", ".xlsx", ".xlsm", ".pptx"}


def load_source_documents(data_dir: Path) -> list[Document]:
	documents: list[Document] = []
	for path in sorted(data_dir.rglob("*")):
		if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
			documents.extend(load_file(path))
	return documents


def load_file(path: Path) -> list[Document]:
	extension = path.suffix.lower()
	if extension == ".pdf":
		return _load_pdf(path)
	if extension == ".docx":
		documents = Docx2txtLoader(str(path)).load()
		return [_add_common_metadata(document, path, "docx") for document in documents]
	if extension in {".txt", ".md"}:
		documents = TextLoader(str(path), encoding="utf-8").load()
		return [_add_common_metadata(document, path, extension[1:]) for document in documents]
	if extension == ".csv":
		return _load_csv(path)
	if extension in {".xlsx", ".xlsm"}:
		return _load_workbook(path)
	if extension == ".pptx":
		return _load_presentation(path)
	raise ValueError(f"Unsupported document extension: {extension}")


def _add_common_metadata(document: Document, path: Path, document_type: str) -> Document:
	document.metadata.update({
		"source": path.name,
		"source_path": path.name,
		"document_type": document_type,
	})
	return document


def _load_pdf(path: Path) -> list[Document]:
	"""Load readable PDF pages independently so one malformed page does not stop ingestion."""
	try:
		reader = PdfReader(str(path), strict=False)
	except Exception as error:
		print(f"Warning: skipped unreadable PDF {path.name}: {error}")
		return []
	try:
		pages = list(reader.pages)
	except Exception as error:
		print(f"Warning: skipped unreadable PDF {path.name}: {error}")
		return []
	documents: list[Document] = []
	for page_number, page in enumerate(pages, start=1):
		try:
			text = (page.extract_text() or "").strip()
		except Exception as error:
			print(f"Warning: skipped unreadable PDF page {page_number} in {path.name}: {error}")
			continue
		if text:
			documents.append(Document(
				page_content=text,
				metadata={
					"source": path.name,
					"source_path": path.name,
					"document_type": "pdf",
					"page_number": page_number,
					"chunk_strategy": "pdf_page",
				},
			))
	if not documents:
		print(f"Warning: no readable text pages found in PDF {path.name}.")
	return documents


def _load_workbook(path: Path) -> list[Document]:
	workbook = load_workbook(path, read_only=True, data_only=True)
	documents: list[Document] = []
	for worksheet in workbook.worksheets:
		rows = worksheet.iter_rows(values_only=True)
		headers = [str(value).strip() if value is not None else f"column_{index + 1}" for index, value in enumerate(next(rows, ())) ]
		for row_number, row in enumerate(rows, start=2):
			values = ["" if value is None else str(value).strip() for value in row]
			if not any(values):
				continue
			fields = [f"{header}: {value}" for header, value in zip(headers, values) if value]
			documents.append(Document(
				page_content=" | ".join(fields),
				metadata={
					"source": path.name,
					"source_path": path.name,
					"document_type": "spreadsheet",
					"sheet_name": worksheet.title,
					"row_number": row_number,
					"chunk_strategy": "excel_row",
				},
			))
	workbook.close()
	return documents


def _load_presentation(path: Path) -> list[Document]:
	presentation = Presentation(str(path))
	documents: list[Document] = []
	for slide_number, slide in enumerate(presentation.slides, start=1):
		text_parts: list[str] = []
		slide_title = slide.shapes.title.text.strip() if slide.shapes.title else ""
		for shape in slide.shapes:
			if shape.has_table:
				for row in shape.table.rows:
					cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
					if cells:
						text_parts.append(" | ".join(cells))
			elif shape.has_text_frame:
				text = shape.text.strip()
				if text and text not in text_parts:
					text_parts.append(text)
		if text_parts:
			documents.append(Document(
				page_content="\n".join(text_parts),
				metadata={
					"source": path.name,
					"source_path": path.name,
					"document_type": "presentation",
					"slide_number": slide_number,
					"slide_title": slide_title,
					"chunk_strategy": "pptx_slide",
				},
			))
	return documents


def _load_csv(path: Path) -> list[Document]:
	"""Load one searchable document per CSV row while preserving column names."""
	documents: list[Document] = []
	with path.open("r", encoding="utf-8-sig", newline="") as csv_file:
		reader = csv.DictReader(csv_file)
		if not reader.fieldnames:
			return documents
		for row_number, row in enumerate(reader, start=2):
			fields = [
				f"{header.strip()}: {value.strip()}"
				for header, value in row.items()
				if header and value and value.strip()
			]
			if not fields:
				continue
			documents.append(Document(
				page_content=" | ".join(fields),
				metadata={
					"source": path.name,
					"source_path": path.name,
					"document_type": "spreadsheet",
					"row_number": row_number,
					"chunk_strategy": "csv_row",
				},
			))
	return documents